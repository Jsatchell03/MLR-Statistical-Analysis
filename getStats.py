# STAGING AREA FOR STATMONKEY

from lxml import etree
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pathlib import Path
import argparse
import statistics
import matplotlib.image as mpimg
import matplotlib.ticker as tck
from operator import itemgetter
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches

# thicken arrows and totals
# Attacking kicks from half to try zone

# Slide show is 16 by 9
# Save all graphs show is optional
# Make sure user know to exit graph images to continue
# Add seperate show fuction show(pathToStat)
# Logic to show using path
# img = mpimg.imread('your_figure.png')

# # Display the image
# plt.imshow(img)

fieldLength = 140
halfFieldLength = 120
fieldWidth = 68
tryZone = 20
halfwayLine = fieldLength / 2
figWidth = 11
figHeight = 6
linebreakKeyPlayers = []
arrowWidth = 1.5
# 2(x - 70)
# Depreciated
def getKickStats(xml_files, teamName):
    playerDict = {}
    for xml_file in xml_files:
        try:
            tree = etree.parse(str(xml_file))
            root = tree.getroot()
            kicks = root.xpath(f"//instance[code='{teamName} Kick']")
            for kick in kicks:
                player = str(
                    kick.xpath("label[group='Player'][position()=1]/text/text()")[0]
                )
                meters = int(
                    kick.xpath("label[group='Kick Metres'][position()=1]/text/text()")[
                        0
                    ]
                )
                if player in playerDict:
                    playerDict[player] = playerDict[player] + 1
                else:
                    playerDict[player] = 1

        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xml_file.name}: {e}")
    print(playerDict)
    
    kickNumber = list(map(lambda x: x[0], playerDict.values()))
    median = statistics.median(kickNumber) 
    
    print(median)
    kickMeters = list(map(lambda x: x[1], playerDict.values()))
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(10, 8))
    fig.suptitle("Player Kicking Statistics", fontsize=14, y=0.97)

    ax1.bar((sorted(playerDict, key=playerDict.get)[::-1]), (sorted(list(map(lambda x: x[0], playerDict.values())))[::-1]))
    ax1.set_ylabel("Number of Kicks")
    ax1.set_title("Total Kicks By Player")
    ax1.tick_params(axis="x", rotation=45)

    ax2.bar((sorted(playerDict, key=playerDict.get)[::-1]), (sorted(list(map(lambda x: x[1], playerDict.values())))[::-1]))
    ax2.set_ylabel("Total Meters")
    ax2.set_title("Total Kick Meters By Player")
    ax2.tick_params(axis="x", rotation=45)

    plt.tight_layout()
    plt.show()
    plt.subplots_adjust(bottom=0.5)

def getLinebreakPhases(xmlFiles, teamName, show=True):
    fig, (ax) = plt.subplots(figsize=(figWidth, figHeight))
    breakPhases = {}
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            linebreaks = root.xpath(
                f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{teamName}' and group='Attacking Quality']]"
            )
            
            for linebreak in linebreaks:
                player = str(
                    linebreak.xpath(
                        "label[group='Player'][position()=1]/text/text()"
                    )[0]
                )
                phase = linebreak.xpath(
                        "label[group='Phase Number'][position()=1]/text/text()"
                    )[0]
                breakPhases[phase] = breakPhases[phase] + 1 if phase in breakPhases else 1

        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")

    plt.title(f"Phase Of Linebreaks")

    breakPhasesKeys = list(breakPhases.keys())
    breakPhasesKeys.sort()
    sortedBreakPhases = {i: breakPhases[i] for i in breakPhasesKeys}

    ax.bar(sortedBreakPhases.keys(), sortedBreakPhases.values())
    ax.set_ylabel("Number of Breaks")
    ax.tick_params(axis="x", rotation=45)
    ax.set_aspect('equal', adjustable='box')
    ax.yaxis.set_major_locator(tck.MultipleLocator())
    
    plt.show()

def getKickPaths(xmlFiles, teamName, show=True):
    fig, ax = plt.subplots(figsize=(figWidth, figHeight))
    drawRugbyPitch(ax, teamName)
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            kicks = root.xpath(f"//instance[code='{teamName} Kick']")
            for kick in kicks:
                xStart = float(kick.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                yStart = fieldWidth - float(kick.xpath("label[group='Y_Start']/text/text()")[0])
                xEnd = float(kick.xpath("label[group='X_End']/text/text()")[0]) + tryZone
                yEnd = fieldWidth - float(kick.xpath("label[group='Y_End']/text/text()")[0])
                descriptor = str(kick.xpath("label[group='Kick Descriptor']/text/text()")[0])
                if descriptor == "Touch Kick":
                    continue
                kickStyle = str(kick.xpath("label[group='Kick Style']/text/text()")[0])
                dx = xEnd - xStart
                dy = yEnd - yStart
                color = ""
                # Windy
                if kickStyle == "Box":
                    color = "#FF85B4"
                else:
                    match descriptor:
                        # Pocket
                        case  "Territorial":
                            color = "#63AAE3"
                        # Ice
                        case "Low":
                            color = "#E15554"
                        # Snow
                        case "Bomb":
                            color = "#E1BC29"
                        # Wedge
                        case "Chip":
                            color = "#2E8A59"
                        # Kick Pass
                        case "Cross Pitch":
                            color = "#7768AE"
                plt.arrow(xStart, yStart, dx, dy, 
                        head_width=2, head_length=1, 
                        fc=color, ec=color, lw=arrowWidth,
                        length_includes_head=True)
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    pocket = mpatches.Patch(color="#63AAE3", label='Pocket')
    windy =  mpatches.Patch(color="#FF85B4", label='Windy')
    ice = mpatches.Patch(color="#E15554", label='Ice')
    snow = mpatches.Patch(color="#E1BC29", label='Snow')
    wedge = mpatches.Patch(color="#2E8A59", label='Wedge')
    kp = mpatches.Patch(color="#7768AE", label='Kick Pass')
    plt.legend(handles=[pocket, windy, ice ,snow, wedge, kp],loc="lower left")
    plt.title('Kick Paths')
    path = f"Stat PNGs/{teamName.replace(" ", "_")}_Kick_Paths.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def getAttackingKickPaths(xmlFiles, teamName, show=True):
    fig, ax = plt.subplots(figsize=(figWidth, figHeight))
    drawHalfPitch(ax)
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            kicks = root.xpath(f"//instance[code='{teamName} Kick']")
            for kick in kicks:
                xStart = float(kick.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                yStart = fieldWidth - float(kick.xpath("label[group='Y_Start']/text/text()")[0])
                xEnd = float(kick.xpath("label[group='X_End']/text/text()")[0]) + tryZone
                yEnd = fieldWidth - float(kick.xpath("label[group='Y_End']/text/text()")[0])
                descriptor = str(kick.xpath("label[group='Kick Descriptor']/text/text()")[0])
                if descriptor == "Touch Kick":
                    continue
                if xStart < 70:
                    continue
                xStart = (xStart - 70) * 2
                xEnd = (xEnd - 70) * 2
                kickStyle = str(kick.xpath("label[group='Kick Style']/text/text()")[0])
                dx = xEnd - xStart
                dy = yEnd - yStart
                color = ""
                # Windy
                if kickStyle == "Box":
                    color = "#FF85B4"
                else:
                    match descriptor:
                        # Pocket
                        case  "Territorial":
                            color = "#63AAE3"
                        # Ice
                        case "Low":
                            color = "#E15554"
                        # Snow
                        case "Bomb":
                            color = "#E1BC29"
                        # Wedge
                        case "Chip":
                            color = "#2E8A59"
                        # Kick Pass
                        case "Cross Pitch":
                            color = "#7768AE"
                plt.arrow(xStart, yStart, dx, dy, 
                        head_width=2, head_length=1, 
                        fc=color, ec=color, lw=arrowWidth,
                        length_includes_head=True)
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    pocket = mpatches.Patch(color="#63AAE3", label='Pocket')
    windy =  mpatches.Patch(color="#FF85B4", label='Windy')
    ice = mpatches.Patch(color="#E15554", label='Ice')
    snow = mpatches.Patch(color="#E1BC29", label='Snow')
    wedge = mpatches.Patch(color="#2E8A59", label='Wedge')
    kp = mpatches.Patch(color="#7768AE", label='Kick Pass')
    plt.legend(handles=[pocket, windy, ice ,snow, wedge, kp],loc="lower left")
    plt.title('Attacking Kick Paths')
    path = f"Stat PNGs/{teamName.replace(" ", "_")}_Kick_Paths.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def getPlayerKickPaths(xmlFiles, teamName, player, show=True):
    fig, ax = plt.subplots(figsize=(figWidth, figHeight))
    drawRugbyPitch(ax, teamName)
    totalKicks = 0
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='{player}' and group='Player']]")
            totalKicks += len(kicks)
            for kick in kicks:
                xStart = float(kick.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                yStart = fieldWidth - float(kick.xpath("label[group='Y_Start']/text/text()")[0])
                xEnd = float(kick.xpath("label[group='X_End']/text/text()")[0]) + tryZone
                yEnd = fieldWidth - float(kick.xpath("label[group='Y_End']/text/text()")[0])
                descriptor = str(kick.xpath("label[group='Kick Descriptor']/text/text()")[0])
                if descriptor == "Touch Kick":
                    continue
                kickStyle = str(kick.xpath("label[group='Kick Style']/text/text()")[0])
                dx = xEnd - xStart
                dy = yEnd - yStart
                color = ""
                # Windy
                if kickStyle == "Box":
                    color = "#FF85B4"
                else:
                    match descriptor:
                        # Pocket
                        case  "Territorial":
                            color = "#63AAE3"
                        # Ice
                        case "Low":
                            color = "#E15554"
                        # Snow
                        case "Bomb":
                            color = "#E1BC29"
                        # Wedge
                        case "Chip":
                            color = "#2E8A59"
                        # Kick Pass
                        case "Cross Pitch":
                            color = "#7768AE"
                plt.arrow(xStart, yStart, dx, dy, 
                        head_width=2, head_length=1, 
                        fc=color,lw=arrowWidth, ec=color,
                        length_includes_head=True)
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    pocket = mpatches.Patch(color="#63AAE3", label='Pocket')
    windy =  mpatches.Patch(color="#FF85B4", label='Windy')
    ice = mpatches.Patch(color="#E15554", label='Ice')
    snow = mpatches.Patch(color="#E1BC29", label='Snow')
    wedge = mpatches.Patch(color="#2E8A59", label='Wedge')
    kp = mpatches.Patch(color="#7768AE", label='Kick Pass')
    plt.legend(handles=[pocket, windy, ice ,snow, wedge, kp],loc="lower left")
    # Number of kicks is wrong 
    plt.title(f'{player} Kick Paths ({totalKicks} Total)')
    path = f"Stat PNGs/{player.replace(" ", "_")}_Kick_Paths.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def getGroupKickPaths(xmlFiles, teamName, type, show=True):
    fig, ax = plt.subplots(figsize=(figWidth, figHeight))
    drawRugbyPitch(ax, teamName)
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            if type == "windy":
                kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='Box' and group='Kick Style'] and label[text!='Touch Kick' and group='Kick Descriptor']]")
                color = "#FF85B4"
                title = "Windy/Box"
            else:
                match type:
                        # Pocket
                        case  "pocket":
                            kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='Territorial' and group='Kick Descriptor'] and label[text='Regular' and group='Kick Style']]")
                            color = "#4D9DE0"
                            title = "Pocket/Long"
                        # Ice
                        case "ice":
                            kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='Low' and group='Kick Descriptor']]")
                            color = "#E15554"
                            title = "Ice/Grubber"
                        # Snow
                        case "snow":
                            kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='Bomb' and group='Kick Descriptor']]")
                            color = "#E1BC29"
                            title = "Snow/Up And Under"
                        # Wedge
                        case "wedge":
                            kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='Chip' and group='Kick Descriptor']]")
                            color = "#3BB273"
                            title = "Wedge/Chip"
                        # Kick Pass
                        case "kp":
                            kicks = root.xpath(f"//instance[code='{teamName} Kick' and label[text='Cross Pitch' and group='Kick Descriptor']]")
                            color = "#7768AE"
                            title = "Kick Pass/Cross"
            for kick in kicks:
                xStart = float(kick.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                yStart = fieldWidth - float(kick.xpath("label[group='Y_Start']/text/text()")[0])
                xEnd = float(kick.xpath("label[group='X_End']/text/text()")[0]) + tryZone
                yEnd = fieldWidth - float(kick.xpath("label[group='Y_End']/text/text()")[0])
                dx = xEnd - xStart
                dy = yEnd - yStart
                plt.arrow(xStart, yStart, dx, dy, 
                        head_width=2, head_length=1, lw=1.5,
                        fc=color, ec=color,
                        length_includes_head=True)
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    plt.title(f'{title} Kick Paths ({len(kicks)} Total)')
    path = f"Stat PNGs/{teamName.replace(" ", "_")}_{type.capitalize()}_Kick_Paths.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def drawRugbyPitch(ax, teamName):
    # Draw halfway line
    ax.plot([halfwayLine, halfwayLine], [0, fieldWidth], color='black', linestyle='-')
    # Draw Channels
    ax.plot([halfwayLine - 2.5, halfwayLine + 2.5],[fieldWidth - 5, fieldWidth - 5],color="black", linestyle="-" )
    ax.plot([halfwayLine - 2.5, halfwayLine + 2.5],[fieldWidth - 15, fieldWidth - 15],color="black", linestyle="-" )
    ax.plot([halfwayLine - 2.5, halfwayLine + 2.5],[15, 15],color="black", linestyle="-" )
    ax.plot([halfwayLine - 2.5, halfwayLine + 2.5],[5, 5],color="black", linestyle="-" )

    # Draw 10-meter lines

    ax.plot([halfwayLine + 10, halfwayLine + 10], [5 - 2.5, 5 + 2.5], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [5 - 2.5, 5 + 2.5], color='black', linestyle='-')
    ax.plot([halfwayLine + 10, halfwayLine + 10], [15 - 2.5, 15 + 2.5], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [15 - 2.5, 15 + 2.5], color='black', linestyle='-')

    ax.plot([halfwayLine + 10, halfwayLine + 10], [fieldWidth - (5 - 2.5), fieldWidth - (5 + 2.5)], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [fieldWidth - (5 - 2.5), fieldWidth - (5 + 2.5)], color='black', linestyle='-')
    ax.plot([halfwayLine + 10, halfwayLine + 10], [fieldWidth - (15 - 2.5), fieldWidth - (15 + 2.5)], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [fieldWidth - (15 - 2.5), fieldWidth - (15 + 2.5)], color='black', linestyle='-')


    # Creating mid 3 dashes for 10 meter line
    ax.plot([halfwayLine + 10, halfwayLine + 10], [fieldWidth/2 - 2.5, fieldWidth/2 + 2.5], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [fieldWidth/2 - 2.5, fieldWidth/2 + 2.5], color='black', linestyle='-')

    ax.plot([halfwayLine + 10, halfwayLine + 10], [fieldWidth/2 - 11.83, fieldWidth/2 - 6.83], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [fieldWidth/2 - 11.83, fieldWidth/2 - 6.83], color='black', linestyle='-')

    ax.plot([halfwayLine + 10, halfwayLine + 10], [fieldWidth/2 + 11.83, fieldWidth/2 + 6.83], color='black', linestyle='-')
    ax.plot([halfwayLine - 10, halfwayLine - 10], [fieldWidth/2 + 11.83, fieldWidth/2 + 6.83], color='black', linestyle='-')
    # Draw Chanels
    ax.plot([halfwayLine + 7.5, halfwayLine + 12.5], [5, 5], color='black', linestyle='-')
    ax.plot([halfwayLine + 7.5, halfwayLine + 12.5], [15, 15], color='black', linestyle='-')
    ax.plot([halfwayLine + 7.5, halfwayLine + 12.5], [fieldWidth - 5, fieldWidth - 5], color='black', linestyle='-')
    ax.plot([halfwayLine + 7.5, halfwayLine + 12.5], [fieldWidth - 15, fieldWidth - 15], color='black', linestyle='-')

    ax.plot([halfwayLine - 7.5, halfwayLine - 12.5], [5, 5], color='black', linestyle='-')
    ax.plot([halfwayLine - 7.5, halfwayLine - 12.5], [15, 15], color='black', linestyle='-')
    ax.plot([halfwayLine - 7.5, halfwayLine - 12.5], [fieldWidth - 5, fieldWidth - 5], color='black', linestyle='-')
    ax.plot([halfwayLine - 7.5, halfwayLine - 12.5], [fieldWidth - 15, fieldWidth - 15], color='black', linestyle='-')

    # Draw 22-meter lines
    ax.plot([42, 42 ], [0, fieldWidth], color='black')
    ax.plot([fieldLength - 42, fieldLength - 42], [0, fieldWidth], color='black')
    # Draw Chanels
    ax.plot([42 - 2.5, 42 + 2.5],[fieldWidth - 5, fieldWidth - 5],color="black", linestyle="-" )
    ax.plot([42 - 2.5, 42 + 2.5],[fieldWidth - 15, fieldWidth - 15],color="black", linestyle="-" )
    ax.plot([42 - 2.5, 42 + 2.5],[15, 15],color="black", linestyle="-" )
    ax.plot([42 - 2.5, 42 + 2.5],[5, 5],color="black", linestyle="-" )

    ax.plot([(fieldLength - 42) - 2.5, (fieldLength - 42) + 2.5],[fieldWidth - 5, fieldWidth - 5],color="black", linestyle="-" )
    ax.plot([(fieldLength - 42) - 2.5, (fieldLength - 42) + 2.5],[fieldWidth - 15, fieldWidth - 15],color="black", linestyle="-" )
    ax.plot([(fieldLength - 42) - 2.5, (fieldLength - 42) + 2.5],[15, 15],color="black", linestyle="-" )
    ax.plot([(fieldLength - 42) - 2.5, (fieldLength - 42) + 2.5],[5, 5],color="black", linestyle="-" )
    # Draw goal lines
    ax.plot([tryZone, tryZone], [0, fieldWidth], color='black', linewidth=3)
    ax.plot([fieldLength - tryZone, fieldLength - tryZone], [0, fieldWidth], color='black', linewidth=3)

    # Draw the field rectangle
    ax.plot([0, fieldLength], [0, 0], color='green', linewidth=3)
    ax.plot([0, fieldLength], [fieldWidth, fieldWidth], color='green', linewidth=3)
    ax.plot([0, 0], [0, fieldWidth], color='green', linewidth=3)
    ax.plot([fieldLength, fieldLength], [0, fieldWidth], color='green', linewidth=3)

    ax.text(tryZone/2, fieldWidth/2, teamName,
            rotation=-90,
            ha='center',
            va='center',
            fontsize=12)
    
    # Right try zone
    ax.text(fieldLength - tryZone/2, fieldWidth/2, 'Opposition',
            rotation=90,
            ha='center',
            va='center',
            fontsize=12)

    # Set limits and aspect ratio
    ax.set_xlim(-5, fieldLength + 5)
    ax.set_ylim(-5, fieldWidth + 5)
    ax.set_aspect('equal', adjustable='box')

    # Remove axes ticks and labels
    ax.set_xticks([])
    ax.set_yticks([])
    plt.tight_layout(pad=2.5)

def drawHalfPitch(ax):
    halfFieldLength = 140
    halfTryZone = 40
    # Draw the field rectangle
    ax.plot([0, halfFieldLength], [0, 0], color='green', linewidth=3)
    ax.plot([0, halfFieldLength], [fieldWidth, fieldWidth], color='green', linewidth=3)
    ax.plot([0, 0], [0, fieldWidth], color='green', linewidth=3)
    ax.plot([halfFieldLength, halfFieldLength], [0, fieldWidth], color='green', linewidth=3)
    # Draw 10-meter lines
    ax.plot([20, 20], [5 - 2.5, 5 + 2.5], color='black', linestyle='-')
    ax.plot([20, 20], [15 - 2.5, 15 + 2.5], color='black', linestyle='-')

    ax.plot([20, 20], [fieldWidth - (5 - 2.5), fieldWidth - (5 + 2.5)], color='black', linestyle='-')
    ax.plot([20, 20], [fieldWidth - (15 - 2.5), fieldWidth - (15 + 2.5)], color='black', linestyle='-')


    # Creating mid 3 dashes for 10 meter line
    ax.plot([20, 20], [fieldWidth/2 - 2.5, fieldWidth/2 + 2.5], color='black', linestyle='-')
    ax.plot([20, 20], [fieldWidth/2 - 11.83, fieldWidth/2 - 6.83], color='black', linestyle='-')
    ax.plot([20, 20], [fieldWidth/2 + 11.83, fieldWidth/2 + 6.83], color='black', linestyle='-')
    # Draw Chanels
    ax.plot([10 + 5, 10 + 15], [5, 5], color='black', linestyle='-')
    ax.plot([10 + 5, 10 + 15], [15, 15], color='black', linestyle='-')
    ax.plot([10 + 5, 10 + 15], [fieldWidth - 5, fieldWidth - 5], color='black', linestyle='-')
    ax.plot([10 + 5, 10 + 15], [fieldWidth - 15, fieldWidth - 15], color='black', linestyle='-')

    # Draw 22-meter lines
    ax.plot([56, 56], [0, fieldWidth], color='black')
    # Draw Chanels
    ax.plot([(56) - 5, (56) + 5],[fieldWidth - 5, fieldWidth - 5],color="black", linestyle="-" )
    ax.plot([(56) - 5, (56) + 5],[fieldWidth - 15, fieldWidth - 15],color="black", linestyle="-" )
    ax.plot([(56) - 5, (56) + 5],[15, 15],color="black", linestyle="-" )
    ax.plot([(56) - 5, (56) + 5],[5, 5],color="black", linestyle="-" )
    # Draw goal lines
    ax.plot([100, 100], [0, fieldWidth], color='black', linewidth=3)
    
    # Right try zone
    ax.text(halfFieldLength - halfTryZone/2, fieldWidth/2, 'Opposition',
            rotation=90,
            ha='center',
            va='center',
            fontsize=12)

    # Set limits and aspect ratio
    ax.set_xlim(-5, halfFieldLength + 5)
    ax.set_ylim(-5, fieldWidth + 5)
    ax.set_aspect('equal', adjustable='box')

    # Remove axes ticks and labels
    ax.set_xticks([])
    ax.set_yticks([])
    plt.tight_layout(pad=2.5)

def getLinebreakLocations(xmlFiles, teamName, show=True):
    fig, ax = plt.subplots(figsize=(figWidth, figHeight))
    drawRugbyPitch(ax, teamName)
    xValues = []
    yValues = []
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            linebreaks = root.xpath(
                f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{teamName}' and group='Attacking Quality']]"
            )
            for linebreak in linebreaks:
                xStart = float(linebreak.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                yStart = fieldWidth - float(linebreak.xpath("label[group='Y_Start']/text/text()")[0])
                xValues.append(xStart)
                yValues.append(yStart)       
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    ax.scatter(xValues, yValues) 
    plt.title('Linebreak Locations')
    path = f"Stat PNGs/{teamName.replace(" ", "_")}_Linebreak_Locations.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path
# Update \/ when class with atributes is implemented
def getPlayerLinebreakLocations(xmlFiles, teamName, player, show=True):
    fig, ax = plt.subplots(figsize=(figWidth, figHeight))
    drawRugbyPitch(ax, teamName)
    xValues = []
    yValues = []
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            linebreaks = root.xpath(
                f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{teamName}' and group='Attacking Quality']]"
            )
            for linebreak in linebreaks:
                playerName = str(
                    linebreak.xpath(
                        "label[group='Player'][position()=1]/text/text()"
                    )[0]
                )
                if player == playerName:
                    xStart = float(linebreak.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                    yStart = fieldWidth - float(linebreak.xpath("label[group='Y_Start']/text/text()")[0])
                    xValues.append(xStart)
                    yValues.append(yStart)       
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    ax.scatter(xValues, yValues) 
    plt.title(f'{player} Linebreak Locations')
    path = f"Stat PNGs/{player.replace(" ", "_")}_Linebreak_Locations.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def getLinebreaksByPlayer(xmlFiles, teamName, show=True):
    plt.figure(figsize=(figWidth, figHeight))
    playerBreaks = {}
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            linebreaks = root.xpath(
                f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{teamName}' and group='Attacking Quality']]"
            )
            
            for linebreak in linebreaks:
                player = str(
                    linebreak.xpath(
                        "label[group='Player'][position()=1]/text/text()"
                    )[0]
                )
                
                playerBreaks[player] = playerBreaks[player] + 1 if player in playerBreaks else 1

        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    print(playerBreaks)
    for player in playerBreaks:
        if playerBreaks[player] > statistics.median(playerBreaks.values()):
            linebreakKeyPlayers.append(player)
    plt.title(f"Number Of Linebreaks By Player")
    print(linebreakKeyPlayers)
  
    sortedPLayerBreaks = OrderedDict(sorted(playerBreaks.items(), key=itemgetter(1), reverse=True))
    plt.bar(sortedPLayerBreaks.keys(), sortedPLayerBreaks.values())
    plt.ylabel("Number of Breaks")
    plt.xticks(rotation=45)  
    plt.tight_layout() 
    plt.subplots_adjust(bottom=0.25) 
    plt.gca().yaxis.set_major_locator(tck.MultipleLocator(base=1))  # Using base=1 as an example
    path = f"Stat PNGs/{teamName.replace(" ", "_")}_Linebreaks_By_Player.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def getMaulMap(xmlFiles, teamName, show=True):
    fig, ax = plt.subplots(figsize=(11, 6))
    drawRugbyPitch(ax, teamName)
    xValues = []
    yValues = []
    maulMetersArr = []
    trueMaulMetersArr = []
    colors = []
    for xmlFile in xmlFiles:
        try:
            tree = etree.parse(str(xmlFile))
            root = tree.getroot()
            mauls = root.xpath(
                f"//instance[code='{teamName} Maul']"
            )
            for maul in mauls:
                xStart = float(maul.xpath("label[group='X_Start']/text/text()")[0]) + tryZone
                yStart = fieldWidth - float(maul.xpath("label[group='Y_Start']/text/text()")[0])
                maulOutcome = str(maul.xpath("label[group='Maul Breakdown Outcome']/text/text()")[0])
                if maulOutcome == "Try Scored":
                    maulMetersArr.append(int(999))
                else:
                    maulMetersArr.append(int(maul.xpath("label[group='Maul Metres']/text/text()")[0]))
                trueMaulMetersArr.append(int(maul.xpath("label[group='Maul Metres']/text/text()")[0]))
                xValues.append(xStart)
                yValues.append(yStart)       
        except etree.XMLSyntaxError as e:
            print(f"Error parsing {xmlFile.name}: {e}")
    avg = (sum(trueMaulMetersArr)) / (len(trueMaulMetersArr))
    for dist in maulMetersArr:
        if dist < avg:
            colors.append("#E15554")
        else:
            colors.append("#3BB273")
    ax.scatter(xValues, yValues, c= colors) 
    pos = mpatches.Patch(color="#3BB273", label=f'> {round(avg,1)} Meters Made/Try Scored')
    neg = mpatches.Patch(color="#E15554", label=f'< {round(avg,1)} Meters Made')
    plt.legend(handles=[pos, neg],loc="lower left")
    plt.title(f'Maul Locations ({round(avg, 1)} Meters Per Maul)')
    path = f"Stat PNGs/{teamName.replace(" ", "_")}_Mauls.png"
    if show:
        plt.show()
    else:
        plt.savefig(path)
        return path

def getGoldZoneReport(xmlFiles, teamName):
    pass

def getScrumStats(xmlFiles, teamName):
    # 2 pie charts for Atacking vs Defensive scrums. Won outright, Penalty won, lost outright, penalty conceded
    pass

def getPenaltyStats(xmlFiles, teamName):
    pass

def addStatToPres(prs, statImgPath, teamName):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    slide.shapes.add_picture("assets/bg.png", left, top, width, height)

    top = Inches(1.5)
    left  = Inches(2.5)
    slide.shapes.add_picture(statImgPath, left, top)

    top = Inches(0)
    left  = Inches(0)
    width = height = Inches(1.5)
    slide.shapes.add_picture("assets/HoundsBadge_LightOnDarkBG.png", left, top, width, height)

    top = Inches(7.25)
    left  = Inches(14.5)
    width = height = Inches(1.5)
    slide.shapes.add_picture("assets/HoundsShield_LightOnDarkBG.png", left, top, width, height)
    
    prs.save(f"{teamName}.pptx")

def main():
    parser = argparse.ArgumentParser(description="Process XML files in a directory")
    parser.add_argument("folder", help="Path to folder containing XML files")
    parser.add_argument(
        "team",
        help="Team name spelt and capitalize the exact way it is referenced in Oval Insights XML",
    )
    
    args = parser.parse_args()

    xml_dir = Path(args.folder)
    xml_files = list(xml_dir.glob("*.xml"))
    trackedTeam = str(args.team)


    # getKickPaths(xml_files, trackedTeam)
    # getGroupKickPaths(xml_files, trackedTeam, "Low")
    # getGroupKickPaths(xml_files, trackedTeam, "Bomb")
    # getGroupKickPaths(xml_files, trackedTeam, "Territorial")
    # getGroupKickPaths(xml_files, trackedTeam, "Chip")
    # getGroupKickPaths(xml_files, trackedTeam, "Cross Pitch")
    # getKickStats(xml_files, trackedTeam)
    # getLinebreakLocations(xml_files, trackedTeam)
    # getMaulMap(xml_files, trackedTeam, show=False)
    # prs = Presentation()
    # addStatToPres(prs, getMaulMap(xml_files, trackedTeam, show=False), trackedTeam)
    # addStatToPres(prs, getKickPaths(xml_files, trackedTeam, show=False), trackedTeam)
    # addStatToPres(prs, getGroupKickPaths(xml_files, trackedTeam, "pocket", show=False), trackedTeam)
    # addStatToPres(prs, getGroupKickPaths(xml_files, trackedTeam, "windy", show=False), trackedTeam)
    # addStatToPres(prs, getGroupKickPaths(xml_files, trackedTeam, "ice", show=False), trackedTeam)
    # addStatToPres(prs, getGroupKickPaths(xml_files, trackedTeam, "snow", show=False), trackedTeam)
    # addStatToPres(prs, getGroupKickPaths(xml_files, trackedTeam, "wedge", show=False), trackedTeam)
    # addStatToPres(prs, getGroupKickPaths(xml_files, trackedTeam, "kp", show=False), trackedTeam)
    # addStatToPres(prs, getPlayerKickPaths(xml_files, trackedTeam, "Jason Robertson", show=False), trackedTeam)
    # addStatToPres(prs, getLinebreakLocations(xml_files, trackedTeam, show=False), trackedTeam)
    # addStatToPres(prs, getLinebreaksByPlayer(xml_files, trackedTeam, show=False), trackedTeam)
    # for player in linebreakKeyPlayers:      
    #     addStatToPres(prs, getPlayerLinebreakLocations(xml_files, trackedTeam, player , show=False), trackedTeam)

    getAttackingKickPaths(xml_files, trackedTeam)



if __name__ == "__main__":
    main()
