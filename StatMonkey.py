from lxml import etree
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pathlib import Path
import argparse
import statistics
import matplotlib.image as mpimg
import matplotlib.ticker as tck
from operator import itemgetter
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches
from Database.MongoDB import Mongo
import logging


class StatMonkey:
    fieldLength = 140
    fieldWidth = 68
    tryZone = 20
    halfwayLine = fieldLength / 2
    figWidth = 11
    figHeight = 6
    arrowWidth = 1.5

    def __init__(self, xmlFiles, teamName, mode="presentation"):
        self.linebreakKeyPlayers = []
        self.mainKickers = []
        self.penalizedProps = []
        self.topTurnovers = []
        self.xmlFiles = xmlFiles
        self.teamName = teamName
        self.mode = mode
        logging.basicConfig(
            level=logging.INFO,
            format="%(asctime)s %(message)s",
        )
        self.logger = logging.getLogger()
        if mode == "presentation":
            self.prs = Presentation()

    def show(self, statPath):
        plt.figure(figsize=(self.figWidth, self.figHeight))
        img = mpimg.imread(statPath)
        plt.imshow(img)
        plt.axis("off")
        plt.tight_layout()
        plt.show()
        plt.close()

    def getAllStats(self):
        results = []
        if self.mode == "presentation":
            results.append(self.getKickStats())
            results.append(self.getKickPaths())
            results.append(self.getAttackingKickPaths())
            results.append(self.getGroupKickPaths("pocket"))
            results.append(self.getGroupKickPaths("windy"))
            results.append(self.getGroupKickPaths("ice"))
            results.append(self.getGroupKickPaths("snow"))
            results.append(self.getGroupKickPaths("wedge"))
            results.append(self.getGroupKickPaths("kp"))
            for player in self.mainKickers:
                results.append(self.getPlayerKickPaths(player))
            results.append(self.get22Stats())
            results.append(self.getLinebreakCountByPlayer())
            results.append(self.getLinebreakPhases())
            results.append(self.getLinebreakLocations())
            for player in self.linebreakKeyPlayers:
                results.append(self.getLinebreakLocationsByPlayer(player))
            results.append(self.getMaulMap())
            results.append(self.getScrumStats())
            results.append(self.getScrumConPens())
            results.append(self.getScrumWonPens())
            for player in self.penalizedProps:
                results.append(self.getScrumPensByPlayer(player))
            results.append(self.getTopTryScorers())
            results.append(self.getTopDefendersBeaten())
            results.append(self.getTopTacklers())
            results.append(self.getTopDomTacklers())
            results.append(self.getTopAssisters())
            results.append(self.getTopCarriers())
            for player in self.topCarriers:
                results.append(self.getCarryBreakdown(player))
            results.append(self.getPlayerTurnoverCount())
            for player in self.topTurnovers[:3]:
                results.append(self.getPlayerTurnoverBD(player))
        return results

    def getKickStats(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Kick_Count_By_Player.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))

        playerKicks = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                kicks = root.xpath(f"//instance[code='{self.teamName} Kick']")
                for kick in kicks:
                    descriptor = str(
                        kick.xpath("label[group='Kick Descriptor']/text/text()")[0]
                    )
                    if descriptor == "Touch Kick":
                        continue
                    player = str(
                        kick.xpath("label[group='Player'][position()=1]/text/text()")[0]
                    )

                    if player in playerKicks:
                        playerKicks[player] = playerKicks[player] + 1
                    else:
                        playerKicks[player] = 1

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        median = statistics.median(playerKicks.values())
        for player in playerKicks:
            if playerKicks[player] > median:
                self.mainKickers.append(player)
        sortedPlayerKicks = OrderedDict(
            sorted(playerKicks.items(), key=itemgetter(1), reverse=True)
        )
        plt.title(f"Number Of Kicks By Player")
        bars = plt.bar(sortedPlayerKicks.keys(), sortedPlayerKicks.values())
        for bar in bars:
            height = bar.get_height()
            if height < 2:
                # Place text above the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height + 0.1,  # Add a small offset above the bar
                    f"{int(height)}",
                    ha="center",
                    va="bottom",  # Align to bottom of text
                    fontweight="bold",
                )
            else:
                # Keep current positioning inside the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    (height / 2),
                    f"{int(height)}",
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
        plt.ylabel("Number of Kicks")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        if self.mode == "database":
            return
        else:
            return path

    def addAllStatsToPres(self, statPathArray):
        for stat in statPathArray:
            self.addStatToPres(stat)

    def getLinebreakLocations(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Linebreak_Locations.png"
        self.logger.info(f"Started {path}")
        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawRugbyPitch(ax)
        xValues = []
        yValues = []
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                linebreaks = root.xpath(
                    f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{self.teamName}' and group='Attacking Quality']]"
                )
                total += len(linebreaks)
                for linebreak in linebreaks:
                    xStart = (
                        float(linebreak.xpath("label[group='X_Start']/text/text()")[0])
                        + self.tryZone
                    )
                    yStart = self.fieldWidth - float(
                        linebreak.xpath("label[group='Y_Start']/text/text()")[0]
                    )
                    xValues.append(xStart)
                    yValues.append(yStart)
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        ax.scatter(xValues, yValues)
        plt.title(f"Linebreak Locations ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getLinebreakLocationsByPlayer(self, player):
        path = f"Stat PNGs/{player.replace(" ", "_")}_Linebreak_Locations.png"
        self.logger.info(f"Started {path}")

        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawRugbyPitch(ax)
        xValues = []
        yValues = []
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                linebreaks = root.xpath(
                    f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{self.teamName}' and group='Attacking Quality']]"
                )
                for linebreak in linebreaks:
                    playerName = str(
                        linebreak.xpath(
                            "label[group='Player'][position()=1]/text/text()"
                        )[0]
                    )
                    if player == playerName:
                        total += 1
                        xStart = (
                            float(
                                linebreak.xpath("label[group='X_Start']/text/text()")[0]
                            )
                            + self.tryZone
                        )
                        yStart = self.fieldWidth - float(
                            linebreak.xpath("label[group='Y_Start']/text/text()")[0]
                        )
                        xValues.append(xStart)
                        yValues.append(yStart)
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        ax.scatter(xValues, yValues)
        plt.title(f"{player} Linebreak Locations")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def get22Stats(self):
        path = f"Stat PNGs/{self.teamName.replace(' ', '_')}_22_Stats.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        scoredEntries = 0
        totalEntries = 0
        totalTrys = 0
        totalPens = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                entries = root.xpath(
                    f"//instance[code='{self.teamName} 22 Entry' and label[text='New Entry' and group='22 Entry']]"
                )
                totalEntries += len(entries)
                trys = root.xpath(f"//instance[code='{self.teamName} Try']")
                scoredEntries += len(trys)
                totalTrys += len(trys)
                penKicks = root.xpath(
                    f"//instance[code='{self.teamName} Goal Kick' and label[text='Penalty Goal' and group='Goal Type'] and label[text='Goal Kicked' and group='Goal Outcome']]"
                )
                scoredEntries += len(penKicks)
                totalPens += len(penKicks)
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        pointsPerEntry = round((((totalTrys * 5) + (3 * totalPens)) / totalEntries), 2)
        plt.pie(
            [totalTrys, totalPens, totalEntries - (totalPens + totalTrys)],
            labels=["Try Scored", "Converted Penalty Kick", "No Points Scored"],
            autopct=lambda p: f"{int(p*sum([totalTrys, totalPens, totalEntries - (totalPens + totalTrys)])/100)}",
        )
        plt.title(f"Gold Zone Efficiency ({pointsPerEntry} Points Per Entry)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def drawRugbyPitch(self, ax):
        self.logger.info(f"Started Drawing Full Pitch")

        # Draw halfway line
        ax.plot(
            [self.halfwayLine, self.halfwayLine],
            [0, self.fieldWidth],
            color="black",
            linestyle="-",
        )
        # Draw Channels
        ax.plot(
            [self.halfwayLine - 2.5, self.halfwayLine + 2.5],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 2.5, self.halfwayLine + 2.5],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 2.5, self.halfwayLine + 2.5],
            [15, 15],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 2.5, self.halfwayLine + 2.5],
            [5, 5],
            color="black",
            linestyle="-",
        )

        # Draw 10-meter lines

        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [5 - 2.5, 5 + 2.5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [5 - 2.5, 5 + 2.5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [15 - 2.5, 15 + 2.5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [15 - 2.5, 15 + 2.5],
            color="black",
            linestyle="-",
        )

        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [self.fieldWidth - (5 - 2.5), self.fieldWidth - (5 + 2.5)],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [self.fieldWidth - (5 - 2.5), self.fieldWidth - (5 + 2.5)],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [self.fieldWidth - (15 - 2.5), self.fieldWidth - (15 + 2.5)],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [self.fieldWidth - (15 - 2.5), self.fieldWidth - (15 + 2.5)],
            color="black",
            linestyle="-",
        )

        # Creating mid 3 dashes for 10 meter line
        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [self.fieldWidth / 2 - 2.5, self.fieldWidth / 2 + 2.5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [self.fieldWidth / 2 - 2.5, self.fieldWidth / 2 + 2.5],
            color="black",
            linestyle="-",
        )

        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [self.fieldWidth / 2 - 11.83, self.fieldWidth / 2 - 6.83],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [self.fieldWidth / 2 - 11.83, self.fieldWidth / 2 - 6.83],
            color="black",
            linestyle="-",
        )

        ax.plot(
            [self.halfwayLine + 10, self.halfwayLine + 10],
            [self.fieldWidth / 2 + 11.83, self.fieldWidth / 2 + 6.83],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 10, self.halfwayLine - 10],
            [self.fieldWidth / 2 + 11.83, self.fieldWidth / 2 + 6.83],
            color="black",
            linestyle="-",
        )
        # Draw Chanels
        ax.plot(
            [self.halfwayLine + 7.5, self.halfwayLine + 12.5],
            [5, 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine + 7.5, self.halfwayLine + 12.5],
            [15, 15],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine + 7.5, self.halfwayLine + 12.5],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine + 7.5, self.halfwayLine + 12.5],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )

        ax.plot(
            [self.halfwayLine - 7.5, self.halfwayLine - 12.5],
            [5, 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 7.5, self.halfwayLine - 12.5],
            [15, 15],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 7.5, self.halfwayLine - 12.5],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [self.halfwayLine - 7.5, self.halfwayLine - 12.5],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )

        # Draw 22-meter lines
        ax.plot([42, 42], [0, self.fieldWidth], color="black")
        ax.plot(
            [self.fieldLength - 42, self.fieldLength - 42],
            [0, self.fieldWidth],
            color="black",
        )
        # Draw Chanels
        ax.plot(
            [42 - 2.5, 42 + 2.5],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [42 - 2.5, 42 + 2.5],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )
        ax.plot([42 - 2.5, 42 + 2.5], [15, 15], color="black", linestyle="-")
        ax.plot([42 - 2.5, 42 + 2.5], [5, 5], color="black", linestyle="-")

        ax.plot(
            [(self.fieldLength - 42) - 2.5, (self.fieldLength - 42) + 2.5],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [(self.fieldLength - 42) - 2.5, (self.fieldLength - 42) + 2.5],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [(self.fieldLength - 42) - 2.5, (self.fieldLength - 42) + 2.5],
            [15, 15],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [(self.fieldLength - 42) - 2.5, (self.fieldLength - 42) + 2.5],
            [5, 5],
            color="black",
            linestyle="-",
        )
        # Draw goal lines
        ax.plot(
            [self.tryZone, self.tryZone],
            [0, self.fieldWidth],
            color="black",
            linewidth=3,
        )
        ax.plot(
            [self.fieldLength - self.tryZone, self.fieldLength - self.tryZone],
            [0, self.fieldWidth],
            color="black",
            linewidth=3,
        )

        # Draw the field rectangle
        ax.plot([0, self.fieldLength], [0, 0], color="green", linewidth=3)
        ax.plot(
            [0, self.fieldLength],
            [self.fieldWidth, self.fieldWidth],
            color="green",
            linewidth=3,
        )
        ax.plot([0, 0], [0, self.fieldWidth], color="green", linewidth=3)
        ax.plot(
            [self.fieldLength, self.fieldLength],
            [0, self.fieldWidth],
            color="green",
            linewidth=3,
        )

        ax.text(
            self.tryZone / 2,
            self.fieldWidth / 2,
            self.teamName,
            rotation=-90,
            ha="center",
            va="center",
            fontsize=12,
        )

        # Right try zone
        ax.text(
            self.fieldLength - self.tryZone / 2,
            self.fieldWidth / 2,
            "Opposition",
            rotation=90,
            ha="center",
            va="center",
            fontsize=12,
        )

        # Set limits and aspect ratio
        ax.set_xlim(-5, self.fieldLength + 5)
        ax.set_ylim(-5, self.fieldWidth + 5)
        ax.set_aspect("equal", adjustable="box")

        # Remove axes ticks and labels
        ax.set_xticks([])
        ax.set_yticks([])
        plt.tight_layout(pad=2.5)
        self.logger.info(f"Finished Drawing Full Pitch")

    def drawHalfPitch(self, ax):
        self.logger.info(f"Started Drawing Half Pitch")

        halfFieldLength = 140
        halfTryZone = 40
        # Draw the field rectangle
        ax.plot([0, halfFieldLength], [0, 0], color="green", linewidth=3)
        ax.plot(
            [0, halfFieldLength],
            [self.fieldWidth, self.fieldWidth],
            color="green",
            linewidth=3,
        )
        ax.plot([0, 0], [0, self.fieldWidth], color="green", linewidth=3)
        ax.plot(
            [halfFieldLength, halfFieldLength],
            [0, self.fieldWidth],
            color="green",
            linewidth=3,
        )
        # Draw 10-meter lines
        ax.plot([20, 20], [5 - 2.5, 5 + 2.5], color="black", linestyle="-")
        ax.plot([20, 20], [15 - 2.5, 15 + 2.5], color="black", linestyle="-")

        ax.plot(
            [20, 20],
            [self.fieldWidth - (5 - 2.5), self.fieldWidth - (5 + 2.5)],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [20, 20],
            [self.fieldWidth - (15 - 2.5), self.fieldWidth - (15 + 2.5)],
            color="black",
            linestyle="-",
        )

        # Creating mid 3 dashes for 10 meter line
        ax.plot(
            [20, 20],
            [self.fieldWidth / 2 - 2.5, self.fieldWidth / 2 + 2.5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [20, 20],
            [self.fieldWidth / 2 - 11.83, self.fieldWidth / 2 - 6.83],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [20, 20],
            [self.fieldWidth / 2 + 11.83, self.fieldWidth / 2 + 6.83],
            color="black",
            linestyle="-",
        )
        # Draw Chanels
        ax.plot([10 + 5, 10 + 15], [5, 5], color="black", linestyle="-")
        ax.plot([10 + 5, 10 + 15], [15, 15], color="black", linestyle="-")
        ax.plot(
            [10 + 5, 10 + 15],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [10 + 5, 10 + 15],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )

        # Draw 22-meter lines
        ax.plot([56, 56], [0, self.fieldWidth], color="black")
        # Draw Chanels
        ax.plot(
            [(56) - 5, (56) + 5],
            [self.fieldWidth - 5, self.fieldWidth - 5],
            color="black",
            linestyle="-",
        )
        ax.plot(
            [(56) - 5, (56) + 5],
            [self.fieldWidth - 15, self.fieldWidth - 15],
            color="black",
            linestyle="-",
        )
        ax.plot([(56) - 5, (56) + 5], [15, 15], color="black", linestyle="-")
        ax.plot([(56) - 5, (56) + 5], [5, 5], color="black", linestyle="-")
        # Draw goal lines
        ax.plot([100, 100], [0, self.fieldWidth], color="black", linewidth=3)

        # Right try zone
        ax.text(
            halfFieldLength - halfTryZone / 2,
            self.fieldWidth / 2,
            "Opposition",
            rotation=90,
            ha="center",
            va="center",
            fontsize=12,
        )

        # Set limits and aspect ratio
        ax.set_xlim(-5, halfFieldLength + 5)
        ax.set_ylim(-5, self.fieldWidth + 5)
        ax.set_aspect("equal", adjustable="box")

        # Remove axes ticks and labels
        ax.set_xticks([])
        ax.set_yticks([])
        plt.tight_layout(pad=2.5)
        self.logger.info(f"Finished Drawing Half Pitch")

    def addStatToPres(self, statImgPath):
        self.logger.info(f"Started Adding {statImgPath} To Pres")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        left = top = Inches(0)
        width = self.prs.slide_width
        height = self.prs.slide_height
        slide.shapes.add_picture("assets/bg.png", left, top, width, height)

        top = Inches(1.5)
        left = Inches(2.5)
        slide.shapes.add_picture(statImgPath, left, top)

        top = Inches(0)
        left = Inches(0)
        width = height = Inches(1.5)
        slide.shapes.add_picture(
            "assets/HoundsBadge_LightOnDarkBG.png", left, top, width, height
        )

        if self.teamName != "Chicago Hounds":
            top = Inches(0)
            left = Inches(2)
            width = height = Inches(1)
            slide.shapes.add_picture(
                f"assets/League Logos/{self.teamName.replace(' ', "_")}.png",
                left,
                top,
                width,
                height,
            )

        top = Inches(7.25)
        left = Inches(14.5)
        width = height = Inches(1.5)
        slide.shapes.add_picture(
            "assets/HoundsShield_LightOnDarkBG.png", left, top, width, height
        )

        self.prs.save(f"{self.teamName}.pptx")
        self.logger.info(f"Finished Adding {statImgPath} To Pres")

    def getMaulMap(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Mauls.png"
        self.logger.info(f"Started {path}")
        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawRugbyPitch(ax)
        xValues = []
        yValues = []
        maulMetersArr = []
        trueMaulMetersArr = []
        colors = []
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                mauls = root.xpath(f"//instance[code='{self.teamName} Maul']")
                for maul in mauls:
                    xStart = (
                        float(maul.xpath("label[group='X_Start']/text/text()")[0])
                        + self.tryZone
                    )
                    yStart = self.fieldWidth - float(
                        maul.xpath("label[group='Y_Start']/text/text()")[0]
                    )
                    maulOutcome = str(
                        maul.xpath("label[group='Maul Breakdown Outcome']/text/text()")[
                            0
                        ]
                    )
                    if maulOutcome == "Try Scored":
                        maulMetersArr.append(int(999))
                    else:
                        maulMetersArr.append(
                            int(maul.xpath("label[group='Maul Metres']/text/text()")[0])
                        )
                    trueMaulMetersArr.append(
                        int(maul.xpath("label[group='Maul Metres']/text/text()")[0])
                    )
                    xValues.append(xStart)
                    yValues.append(yStart)
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        avg = (sum(trueMaulMetersArr)) / (len(trueMaulMetersArr))
        for dist in maulMetersArr:
            if dist < avg:
                colors.append("#E15554")
            else:
                colors.append("#3BB273")
        ax.scatter(xValues, yValues, c=colors)
        pos = mpatches.Patch(
            color="#3BB273", label=f"> {round(avg,1)} Meters Made/Try Scored"
        )
        neg = mpatches.Patch(color="#E15554", label=f"< {round(avg,1)} Meters Made")
        plt.legend(handles=[pos, neg], loc="lower left")
        plt.title(f"Maul Locations ({round(avg, 1)} Meters Per Maul)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getLinebreakCountByPlayer(self):

        path = (
            f"Stat PNGs/{self.teamName.replace(" ", "_")}_Linebreak_Count_By_Player.png"
        )
        self.logger.info(f"Started {path}")

        plt.figure(figsize=(self.figWidth, self.figHeight))
        playerBreaks = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                linebreaks = root.xpath(
                    f"""//instance[label[text="Initial Break" and group="Attacking Qualities"] and label[text="{self.teamName}" and group="Attacking Quality"]]"""
                )

                for linebreak in linebreaks:
                    player = str(
                        linebreak.xpath(
                            "label[group='Player'][position()=1]/text/text()"
                        )[0]
                    )

                    playerBreaks[player] = (
                        playerBreaks[player] + 1 if player in playerBreaks else 1
                    )

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        for player in playerBreaks:
            if playerBreaks[player] > statistics.median(playerBreaks.values()):
                self.linebreakKeyPlayers.append(player)
        plt.title(f"Number Of Linebreaks By Player")
        sortedPlayerBreaks = OrderedDict(
            sorted(playerBreaks.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(sortedPlayerBreaks.keys(), sortedPlayerBreaks.values())
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Number of Breaks")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getLinebreakPhases(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Linebreak_Phases.png"
        self.logger.info(f"Started {path}")

        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        breakPhases = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                linebreaks = root.xpath(
                    f"//instance[label[text='Initial Break' and group='Attacking Qualities'] and label[text='{self.teamName}' and group='Attacking Quality']]"
                )

                for linebreak in linebreaks:
                    phase = linebreak.xpath(
                        "label[group='Phase Number'][position()=1]/text/text()"
                    )[0]
                    breakPhases[phase] = (
                        breakPhases[phase] + 1 if phase in breakPhases else 1
                    )

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")

        plt.title(f"Phase Of Linebreaks")

        sortedBreakPhases = OrderedDict(
            sorted(breakPhases.items(), key=itemgetter(1), reverse=True)
        )

        bars = ax.bar(sortedBreakPhases.keys(), sortedBreakPhases.values())
        for bar in bars:
            height = bar.get_height()
            if height < 2:
                # Place text above the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height + 0.1,  # Add a small offset above the bar
                    f"{int(height)}",
                    ha="center",
                    va="bottom",  # Align to bottom of text
                    fontweight="bold",
                )
            else:
                # Keep current positioning inside the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    (height / 2),
                    f"{int(height)}",
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
        ax.set_ylabel("Number of Breaks")
        ax.tick_params(axis="x", rotation=45)

        # ax.set_aspect('equal', adjustable='box')
        ax.yaxis.set_major_locator(tck.MultipleLocator())
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getKickPaths(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Kick_Paths.png"
        self.logger.info(f"Started {path}")

        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawRugbyPitch(ax)
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                kicks = root.xpath(f"//instance[code='{self.teamName} Kick']")
                for kick in kicks:
                    xStart = (
                        float(kick.xpath("label[group='X_Start']/text/text()")[0])
                        + self.tryZone
                    )
                    yStart = self.fieldWidth - float(
                        kick.xpath("label[group='Y_Start']/text/text()")[0]
                    )
                    xEnd = (
                        float(kick.xpath("label[group='X_End']/text/text()")[0])
                        + self.tryZone
                    )
                    yEnd = self.fieldWidth - float(
                        kick.xpath("label[group='Y_End']/text/text()")[0]
                    )
                    descriptor = str(
                        kick.xpath("label[group='Kick Descriptor']/text/text()")[0]
                    )
                    if descriptor == "Touch Kick":
                        continue
                    total += 1
                    kickStyle = str(
                        kick.xpath("label[group='Kick Style']/text/text()")[0]
                    )
                    dx = xEnd - xStart
                    dy = yEnd - yStart
                    color = ""
                    # Windy
                    if kickStyle == "Box":
                        color = "#FF85B4"
                    else:
                        match descriptor:
                            # Pocket
                            case "Territorial":
                                color = "#63AAE3"
                            # Ice
                            case "Low":
                                color = "#E15554"
                            # Snow
                            case "Bomb":
                                color = "#E1BC29"
                            # Wedge
                            case "Chip":
                                color = "#2E8A59"
                            # Kick Pass
                            case "Cross Pitch":
                                color = "#7768AE"
                    plt.arrow(
                        xStart,
                        yStart,
                        dx,
                        dy,
                        head_width=2,
                        head_length=1,
                        fc=color,
                        ec=color,
                        lw=self.arrowWidth,
                        length_includes_head=True,
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        pocket = mpatches.Patch(color="#63AAE3", label="Pocket")
        windy = mpatches.Patch(color="#FF85B4", label="Windy")
        ice = mpatches.Patch(color="#E15554", label="Ice")
        snow = mpatches.Patch(color="#E1BC29", label="Snow")
        wedge = mpatches.Patch(color="#2E8A59", label="Wedge")
        kp = mpatches.Patch(color="#7768AE", label="Kick Pass")
        plt.legend(handles=[pocket, windy, ice, snow, wedge, kp], loc="lower left")
        plt.title(f"Kick Paths ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getAttackingKickPaths(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Attacking_Kick_Paths.png"
        self.logger.info(f"Started {path}")
        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawHalfPitch(ax)
        total = 0

        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                kicks = root.xpath(f"//instance[code='{self.teamName} Kick']")
                for kick in kicks:
                    xStart = (
                        float(kick.xpath("label[group='X_Start']/text/text()")[0])
                        + self.tryZone
                    )
                    yStart = self.fieldWidth - float(
                        kick.xpath("label[group='Y_Start']/text/text()")[0]
                    )
                    xEnd = (
                        float(kick.xpath("label[group='X_End']/text/text()")[0])
                        + self.tryZone
                    )
                    yEnd = self.fieldWidth - float(
                        kick.xpath("label[group='Y_End']/text/text()")[0]
                    )
                    descriptor = str(
                        kick.xpath("label[group='Kick Descriptor']/text/text()")[0]
                    )
                    if descriptor == "Touch Kick":
                        continue
                    if xStart < 70:
                        continue
                    total += 1
                    xStart = (xStart - 70) * 2
                    xEnd = (xEnd - 70) * 2
                    kickStyle = str(
                        kick.xpath("label[group='Kick Style']/text/text()")[0]
                    )
                    dx = xEnd - xStart
                    dy = yEnd - yStart
                    color = ""
                    # Windy
                    if kickStyle == "Box":
                        color = "#FF85B4"

                    else:
                        match descriptor:
                            # Pocket
                            case "Territorial":
                                color = "#63AAE3"
                            # Ice
                            case "Low":
                                color = "#E15554"
                            # Snow
                            case "Bomb":
                                color = "#E1BC29"
                            # Wedge
                            case "Chip":
                                color = "#2E8A59"
                            # Kick Pass
                            case "Cross Pitch":
                                color = "#7768AE"
                    plt.arrow(
                        xStart,
                        yStart,
                        dx,
                        dy,
                        head_width=2,
                        head_length=1,
                        fc=color,
                        ec=color,
                        lw=self.arrowWidth,
                        length_includes_head=True,
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        pocket = mpatches.Patch(color="#63AAE3", label="Pocket")
        windy = mpatches.Patch(color="#FF85B4", label="Windy")
        ice = mpatches.Patch(color="#E15554", label="Ice")
        snow = mpatches.Patch(color="#E1BC29", label="Snow")
        wedge = mpatches.Patch(color="#2E8A59", label="Wedge")
        kp = mpatches.Patch(color="#7768AE", label="Kick Pass")
        plt.legend(handles=[pocket, windy, ice, snow, wedge, kp], loc="lower left")
        plt.title(f"Attacking Kick Paths ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getPlayerKickPaths(self, player):
        path = f"Stat PNGs/{player.replace(" ", "_")}_Kick_Paths.png"
        self.logger.info(f"Started {path}")

        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawRugbyPitch(ax)
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                kicks = root.xpath(
                    f"""//instance[code="{self.teamName} Kick" and label[text="{player}" and group="Player"]]"""
                )
                for kick in kicks:
                    xStart = (
                        float(kick.xpath("label[group='X_Start']/text/text()")[0])
                        + self.tryZone
                    )
                    yStart = self.fieldWidth - float(
                        kick.xpath("label[group='Y_Start']/text/text()")[0]
                    )
                    xEnd = (
                        float(kick.xpath("label[group='X_End']/text/text()")[0])
                        + self.tryZone
                    )
                    yEnd = self.fieldWidth - float(
                        kick.xpath("label[group='Y_End']/text/text()")[0]
                    )
                    descriptor = str(
                        kick.xpath("label[group='Kick Descriptor']/text/text()")[0]
                    )
                    if descriptor == "Touch Kick":
                        continue
                    total += 1
                    kickStyle = str(
                        kick.xpath("label[group='Kick Style']/text/text()")[0]
                    )
                    dx = xEnd - xStart
                    dy = yEnd - yStart
                    # Windy
                    if kickStyle == "Box":
                        color = "#FF85B4"
                    else:
                        match descriptor:
                            # Pocket
                            case "Territorial":
                                color = "#63AAE3"
                            # Ice
                            case "Low":
                                color = "#E15554"
                            # Snow
                            case "Bomb":
                                color = "#E1BC29"
                            # Wedge
                            case "Chip":
                                color = "#2E8A59"
                            # Kick Pass
                            case "Cross Pitch":
                                color = "#7768AE"
                    plt.arrow(
                        xStart,
                        yStart,
                        dx,
                        dy,
                        head_width=2,
                        head_length=1,
                        fc=color,
                        ec=color,
                        lw=self.arrowWidth,
                        length_includes_head=True,
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        pocket = mpatches.Patch(color="#63AAE3", label="Pocket")
        windy = mpatches.Patch(color="#FF85B4", label="Windy")
        ice = mpatches.Patch(color="#E15554", label="Ice")
        snow = mpatches.Patch(color="#E1BC29", label="Snow")
        wedge = mpatches.Patch(color="#2E8A59", label="Wedge")
        kp = mpatches.Patch(color="#7768AE", label="Kick Pass")
        plt.legend(handles=[pocket, windy, ice, snow, wedge, kp], loc="lower left")
        plt.title(f"{player} Kick Paths ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getGroupKickPaths(self, type):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_{type.capitalize()}_Kick_Paths.png"
        self.logger.info(f"Started {path}")

        fig, ax = plt.subplots(figsize=(self.figWidth, self.figHeight))
        self.drawRugbyPitch(ax)
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                if type == "windy":
                    kicks = root.xpath(
                        f"//instance[code='{self.teamName} Kick' and label[text='Box' and group='Kick Style'] and label[text!='Touch Kick' and group='Kick Descriptor']]"
                    )
                    color = "#FF85B4"
                    title = "Windy/Box"
                else:
                    match type:
                        # Pocket
                        case "pocket":
                            kicks = root.xpath(
                                f"//instance[code='{self.teamName} Kick' and label[text='Territorial' and group='Kick Descriptor'] and label[text='Regular' and group='Kick Style']]"
                            )
                            color = "#4D9DE0"
                            title = "Pocket/Long"
                        # Ice
                        case "ice":
                            kicks = root.xpath(
                                f"//instance[code='{self.teamName} Kick' and label[text='Low' and group='Kick Descriptor']]"
                            )
                            color = "#E15554"
                            title = "Ice/Grubber"
                        # Snow
                        case "snow":
                            kicks = root.xpath(
                                f"//instance[code='{self.teamName} Kick' and label[text='Bomb' and group='Kick Descriptor']]"
                            )
                            color = "#E1BC29"
                            title = "Snow/Up And Under"
                        # Wedge
                        case "wedge":
                            kicks = root.xpath(
                                f"//instance[code='{self.teamName} Kick' and label[text='Chip' and group='Kick Descriptor']]"
                            )
                            color = "#3BB273"
                            title = "Wedge/Chip"
                        # Kick Pass
                        case "kp":
                            kicks = root.xpath(
                                f"//instance[code='{self.teamName} Kick' and label[text='Cross Pitch' and group='Kick Descriptor']]"
                            )
                            color = "#7768AE"
                            title = "Kick Pass/Cross"
                for kick in kicks:
                    total += 1
                    xStart = (
                        float(kick.xpath("label[group='X_Start']/text/text()")[0])
                        + self.tryZone
                    )
                    yStart = self.fieldWidth - float(
                        kick.xpath("label[group='Y_Start']/text/text()")[0]
                    )
                    xEnd = (
                        float(kick.xpath("label[group='X_End']/text/text()")[0])
                        + self.tryZone
                    )
                    yEnd = self.fieldWidth - float(
                        kick.xpath("label[group='Y_End']/text/text()")[0]
                    )
                    dx = xEnd - xStart
                    dy = yEnd - yStart
                    plt.arrow(
                        xStart,
                        yStart,
                        dx,
                        dy,
                        head_width=2,
                        head_length=1,
                        fc=color,
                        ec=color,
                        lw=self.arrowWidth,
                        length_includes_head=True,
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.title(f"{title} Kick Paths ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getScrumStats(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Scrum_Stats.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))

        scrumStats = {
            "Won Outright": 0,
            "Won Penalty": 0,
            "Lost Outright": 0,
            "Conceded Penalty": 0,
            "Reset": 0,
        }
        totalScrums = 0
        positiveScrums = 0
        negativeScrums = 0

        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                scrums = root.xpath(f"//instance[code='{self.teamName} Scrum']")
                for scrum in scrums:
                    result = str(
                        scrum.xpath("label[group='Scrum Result']/text/text()")[0]
                    )
                    totalScrums += 1
                    match result:
                        case "Reset":
                            scrumStats["Reset"] = scrumStats["Reset"] + 1
                        case "Won Outright" | "Won Try":
                            scrumStats["Won Outright"] = scrumStats["Won Outright"] + 1
                            positiveScrums += 1
                        case "Won Free Kick" | "Won Penalty" | "Won Penalty Try":
                            scrumStats["Won Penalty"] = scrumStats["Won Penalty"] + 1
                            positiveScrums += 1
                        case "Lost Outright":
                            scrumStats["Lost Outright"] = (
                                scrumStats["Lost Outright"] + 1
                            )
                            negativeScrums += 1
                        case "Lost Pen Con" | "Lost Free Kick":
                            scrumStats["Conceded Penalty"] = (
                                scrumStats["Conceded Penalty"] + 1
                            )
                            negativeScrums += 1
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        successRate = int(round(positiveScrums / totalScrums, 2) * 100)
        sortedScrumStats = OrderedDict(
            sorted(scrumStats.items(), key=itemgetter(1), reverse=True)
        )
        plt.title(
            f"{self.teamName} Attacking Scrum Results ({successRate}% Success {positiveScrums}/{totalScrums})"
        )
        bars = plt.bar(sortedScrumStats.keys(), sortedScrumStats.values())
        for bar in bars:
            height = bar.get_height()
            if height < 2:
                # Place text above the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height + 0.1,  # Add a small offset above the bar
                    f"{int(height)}",
                    ha="center",
                    va="bottom",  # Align to bottom of text
                    fontweight="bold",
                )
            else:
                # Keep current positioning inside the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    (height / 2),
                    f"{int(height)}",
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getScrumConPens(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Conceded_Scrum_Pens.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        penaltyCount = {}
        totalPens = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                pens = root.xpath(
                    f"//instance[code='{self.teamName} Penalty Conceded' and label[text='Scrum Offence' and group='Pen Descriptor']]"
                )
                for pen in pens:
                    offence = str(
                        pen.xpath("label[group='Scrum Offences']/text/text()")[0]
                    )
                    if pen.xpath("label[group='Player']/text/text()"):
                        player = str(pen.xpath("label[group='Player']/text/text()")[0])
                        if player not in self.penalizedProps:
                            self.penalizedProps.append(player)

                    totalPens += 1
                    if offence not in penaltyCount:
                        penaltyCount[offence] = 1
                    else:
                        penaltyCount[offence] = penaltyCount[offence] + 1

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        sortedPenCount = OrderedDict(
            sorted(penaltyCount.items(), key=itemgetter(1), reverse=True)
        )
        plt.title(f"{self.teamName} Scrum Penalties Conceded ({totalPens} Total)")
        bars = plt.bar(sortedPenCount.keys(), sortedPenCount.values())
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.xticks(rotation=45)
        plt.gca().yaxis.set_major_locator(plt.MultipleLocator(1))
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getScrumWonPens(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Scrum_Pens_Won.png"
        self.logger.info(f"Started {path}")

        plt.figure(figsize=(self.figWidth, self.figHeight))
        penaltyCount = {}
        totalPens = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                opp = root.xpath(
                    f"string(//label/text[following-sibling::group[text()='Penalty Conceded' or text()='Turnover Won' or text()='Kick' or text()='Pass' or text()='Restart']][not(text()='{self.teamName}')][1])"
                )
                pens = root.xpath(
                    f"//instance[code='{opp} Penalty Conceded' and label[text='Scrum Offence' and group='Pen Descriptor']]"
                )
                for pen in pens:
                    offence = str(
                        pen.xpath("label[group='Scrum Offences']/text/text()")[0]
                    )
                    totalPens += 1
                    if offence not in penaltyCount:
                        penaltyCount[offence] = 1
                    else:
                        penaltyCount[offence] = penaltyCount[offence] + 1
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        sortedPenCount = OrderedDict(
            sorted(penaltyCount.items(), key=itemgetter(1), reverse=True)
        )
        plt.title(f"{self.teamName} Scrum Penalties Won ({totalPens} Total)")
        bars = plt.bar(sortedPenCount.keys(), sortedPenCount.values())
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.xticks(rotation=45)
        plt.gca().yaxis.set_major_locator(plt.MultipleLocator(1))
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getScrumPensByPlayer(self, player):
        path = f"Stat PNGs/{player.replace(" ", "_")}_Scrum_Pens.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        playerPens = {}
        totalPens = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                pens = root.xpath(
                    f"//instance[code='{self.teamName} Penalty Conceded' and label[text='Scrum Offence' and group='Pen Descriptor'] and label[text='{player}' and group='Player']]"
                )
                for pen in pens:
                    offence = str(
                        pen.xpath("label[group='Scrum Offences']/text/text()")[0]
                    )
                    totalPens += 1
                    if offence not in playerPens:
                        playerPens[offence] = 1
                    else:
                        playerPens[offence] = playerPens[offence] + 1

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        sortedPlayerPens = OrderedDict(
            sorted(playerPens.items(), key=itemgetter(1), reverse=True)
        )
        plt.title(f"{player} Scrum Penalties Conceded ({totalPens} Total)")
        plt.pie(
            sortedPlayerPens.values(),
            labels=sortedPlayerPens.keys(),
            autopct=lambda p: f"{int(int(p*sum(sortedPlayerPens.values())) / 100)}",
        )
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getTopDefendersBeaten(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Defenders_Beaten.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        defenderBeaters = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                defenceBeaten = root.xpath(
                    f"//instance[label[text='Defender Beaten' and group='Attacking Qualities'] and label[text='{self.teamName}' and group='Attacking Quality']]"
                )
                for instance in defenceBeaten:
                    player = str(
                        instance.xpath(
                            "label[group='Player'][position()=1]/text/text()"
                        )[0]
                    )
                    defenderBeaters[player] = (
                        defenderBeaters[player] + 1 if player in defenderBeaters else 1
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.title(f"Top Performers: Defenders Beaten")
        sortedDefenderBeaters = OrderedDict(
            sorted(defenderBeaters.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(
            list(sortedDefenderBeaters.keys())[:5],
            list(sortedDefenderBeaters.values())[:5],
        )
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Defenders Beaten")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Defenders_Beaten.png"
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getTopTryScorers(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Try_Scorers.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        tryScorers = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                trys = root.xpath(f"//instance[code='{self.teamName} Try']")

                for instance in trys:
                    player = instance.xpath(
                        "label[group='Player'][position()=1]/text/text()"
                    )
                    if player:
                        player = player[0]
                    else:
                        continue
                    tryScorers[player] = (
                        tryScorers[player] + 1 if player in tryScorers else 1
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.title(f"Top Performers: Try Scorers")
        sortedTryScorers = OrderedDict(
            sorted(tryScorers.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(
            list(sortedTryScorers.keys())[:5],
            list(sortedTryScorers.values())[:5],
        )
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Tries Scored")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.gca().yaxis.set_major_locator(tck.MultipleLocator(base=1))
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getTopTacklers(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Tacklers.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        tacklers = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                tackles = root.xpath(
                    f"//instance[label[text='Complete' and group='Tackle Outcome'] and label[text='{self.teamName}' and group='Tackle'] and label[text='Tackle' and group='Event']]"
                )
                for tackle in tackles:
                    player = str(tackle.xpath("label[group='Player']/text/text()")[0])
                    tacklers[player] = tacklers[player] + 1 if player in tacklers else 1
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.title(f"Top Performers: Completed Tackles")
        sortedTacklers = OrderedDict(
            sorted(tacklers.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(
            list(sortedTacklers.keys())[:5],
            list(sortedTacklers.values())[:5],
        )
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Completed Tackles")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getTopDomTacklers(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Dom_Tacklers.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        tacklers = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                tackles = root.xpath(
                    f"//instance[label[text='Complete' and group='Tackle Outcome'] and label[text='{self.teamName}' and group='Tackle'] and label[text='Tackle' and group='Event'] and label[text='Dominant Tackle Contact' and group='Tackle Dominance']]"
                )
                for tackle in tackles:
                    player = str(tackle.xpath("label[group='Player']/text/text()")[0])
                    tacklers[player] = tacklers[player] + 1 if player in tacklers else 1
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.title(f"Top Performers: Dominant Tackles")
        sortedTacklers = OrderedDict(
            sorted(tacklers.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(
            list(sortedTacklers.keys())[:5],
            list(sortedTacklers.values())[:5],
        )
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Dominant Tackles")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getTopAssisters(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Assisters.png"
        self.logger.info(f"Started {path}")

        plt.figure(figsize=(self.figWidth, self.figHeight))
        assisters = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                assists = root.xpath(
                    f"//instance[label[text='Try Assist' and group='Attacking Qualities'] and label[text='{self.teamName}' and group='Attacking Quality']]"
                )
                for assist in assists:
                    player = str(
                        assist.xpath("label[group='Player'][position()=1]/text/text()")[
                            0
                        ]
                    )
                    assisters[player] = (
                        assisters[player] + 1 if player in assisters else 1
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        if len(assisters.keys()) == 0:
            return None
        plt.title(f"Top Performers: Assists")
        sortedAssisters = OrderedDict(
            sorted(assisters.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(
            list(sortedAssisters.keys())[:5],
            list(sortedAssisters.values())[:5],
        )
        self.topAssisters = list(sortedAssisters.keys())[:5]
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Assists")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.gca().yaxis.set_major_locator(tck.MultipleLocator(base=1))
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    # WIP get more info about the pass
    def getAssistBreakdown(self, player):
        path = f"Stat PNGs/{player.replace(' ', '_')}_Assist_Breakdown.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        assistStyles = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                assists = root.xpath(
                    f"//instance[label[text='Try Assist' and group='Attacking Qualities'] and label[text='{self.teamName}' and group='Attacking Quality'] and label[text='{player}' and group='Player']]"
                )
                for assist in assists:
                    style = str(
                        assist.xpath(
                            "label[group='Assist Style'][position()=1]/text/text()"
                        )[0]
                    )
                    assistStyles[style] = (
                        assistStyles[style] + 1 if style in assistStyles else 1
                    )
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.pie(
            assistStyles.values(),
            labels=assistStyles.keys(),
            autopct=lambda p: f"{int(p*sum(assistStyles.values())) / 100}",
        )

        plt.title(f"{player} Assist Breakdown ({len(assistStyles.keys())} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getTopCarriers(self):
        path = f"Stat PNGs/{self.teamName.replace(" ", "_")}_Top_Carriers.png"
        self.logger.info(f"Started {path}")

        plt.figure(figsize=(self.figWidth, self.figHeight))
        carriers = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                carries = root.xpath(
                    f"//instance[label[text='{self.teamName}' and group='Carry'] and label[text='Carry' and group='Event']]"
                )
                for carry in carries:
                    player = str(carry.xpath("label[group='Player']/text/text()")[0])
                    carriers[player] = carriers[player] + 1 if player in carriers else 1
            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.title(f"Top Performers: Carries")
        sortedCarriers = OrderedDict(
            sorted(carriers.items(), key=itemgetter(1), reverse=True)
        )
        self.topCarriers = list(sortedCarriers.keys())[:5]
        bars = plt.bar(
            list(sortedCarriers.keys())[:5],
            list(sortedCarriers.values())[:5],
        )
        for bar in bars:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2.0,
                (height / 2),
                f"{int(height)}",
                ha="center",
                va="center",
                fontweight="bold",
            )
        plt.ylabel("Carries")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.25)
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")

        return path

    def getCarryBreakdown(self, player):
        path = f"Stat PNGs/{player.replace(' ', '_')}_Carry_Breakdown.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        breakdown = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                carries = root.xpath(
                    f"""//instance[label[text="{self.teamName}" and group="Carry"] and label[text="Carry" and group="Event"] and label[text="{player}" and group="Player"]]"""
                )
                for carry in carries:
                    outcome = str(
                        carry.xpath("label[group='Carry Outcome']/text/text()")[0]
                    )
                    if outcome == "Other":
                        continue
                    elif outcome == "Tackled":
                        contact = str(
                            carry.xpath("label[group='Carry Dominance']/text/text()")[0]
                        )
                        breakdown[contact] = (
                            breakdown[contact] + 1 if contact in breakdown else 1
                        )
                    else:
                        breakdown[outcome] = (
                            breakdown[outcome] + 1 if outcome in breakdown else 1
                        )

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.pie(
            breakdown.values(),
            labels=breakdown.keys(),
            autopct=lambda p: f"{int(int(p*sum(breakdown.values())) / 100)}",
        )

        plt.title(f"{player} Carries Breakdown")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    # WIP
    def getTurnoverStats(self):
        path = f"Stat PNGs/{self.teamName.replace(' ', '_')}_Turnover_Breakdown.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        breakdown = {}
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                turnovers = root.xpath(f"//instance[code='{self.teamName} Turnover']")

                for turnover in turnovers:
                    total += 1
                    descriptor = str(
                        turnover.xpath("label[group='Error Descriptor']/text/text()")[0]
                    )
                    if "Kick" in descriptor:
                        descriptor = "Kick Error"

                    breakdown[descriptor] = (
                        breakdown[descriptor] + 1 if descriptor in breakdown else 1
                    )

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        sortedBreakdown = OrderedDict(
            sorted(breakdown.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(sortedBreakdown.keys(), sortedBreakdown.values())
        for bar in bars:
            height = bar.get_height()
            if height < 2:
                # Place text above the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height + 0.1,  # Add a small offset above the bar
                    f"{int(height)}",
                    ha="center",
                    va="bottom",  # Align to bottom of text
                    fontweight="bold",
                )
            else:
                # Keep current positioning inside the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    (height / 2),
                    f"{int(height)}",
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
        plt.xticks(rotation=45)
        plt.subplots_adjust(bottom=0.25)
        plt.title(f"{self.teamName} Turnover Breakdown ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getPlayerTurnoverCount(self):
        path = f"Stat PNGs/{self.teamName.replace(' ', '_')}_Turnover_Count.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        breakdown = {}
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                turnovers = root.xpath(f"//instance[code='{self.teamName} Turnover']")

                for turnover in turnovers:
                    player = str(turnover.xpath("label[group='Player']/text/text()")[0])

                    breakdown[player] = (
                        breakdown[player] + 1 if player in breakdown else 1
                    )

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        median = statistics.median(breakdown.values())
        for player in breakdown:
            if breakdown[player] > median:
                self.topTurnovers.append(player)
        sortedBreakdown = OrderedDict(
            sorted(breakdown.items(), key=itemgetter(1), reverse=True)
        )
        bars = plt.bar(sortedBreakdown.keys(), sortedBreakdown.values())
        for bar in bars:
            height = bar.get_height()
            if height < 2:
                # Place text above the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height + 0.1,  # Add a small offset above the bar
                    f"{int(height)}",
                    ha="center",
                    va="bottom",  # Align to bottom of text
                    fontweight="bold",
                )
            else:
                # Keep current positioning inside the bar
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    (height / 2),
                    f"{int(height)}",
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
        plt.xticks(rotation=60)
        plt.subplots_adjust(bottom=0.25)
        plt.title(f"{self.teamName} Player Turnover Count")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path

    def getPlayerTurnoverBD(self, player):
        path = f"Stat PNGs/{player}_Turnover_Breakdown.png"
        self.logger.info(f"Started {path}")
        plt.figure(figsize=(self.figWidth, self.figHeight))
        breakdown = {}
        total = 0
        for xmlFile in self.xmlFiles:
            try:
                tree = etree.parse(str(xmlFile))
                root = tree.getroot()
                turnovers = root.xpath(
                    f"""//instance[code="{self.teamName} Turnover" and label[text="{player}" and group="Player"]]"""
                )
                for turnover in turnovers:
                    total += 1
                    descriptor = str(
                        turnover.xpath("label[group='Error Descriptor']/text/text()")[0]
                    )
                    if "Kick" in descriptor:
                        descriptor = "Kick Error"

                    breakdown[descriptor] = (
                        breakdown[descriptor] + 1 if descriptor in breakdown else 1
                    )

            except etree.XMLSyntaxError as e:
                print(f"Error parsing {xmlFile.name}: {e}")
        plt.pie(
            breakdown.values(),
            labels=breakdown.keys(),
            autopct=lambda p: f"{int(int(p*sum(breakdown.values())) / 100)}",
        )
        plt.title(f"{player} Turnover Breakdown ({total} Total)")
        plt.savefig(path)
        plt.close()
        self.logger.info(f"Finished {path}")
        return path


def main():
    parser = argparse.ArgumentParser(description="Process XML files in a directory")
    parser.add_argument("folder", help="Path to folder containing XML files")
    parser.add_argument(
        "team",
        help="Team name spelt and capitalize the exact way it is referenced in Oval Insights XML",
    )

    args = parser.parse_args()

    xml_dir = Path(args.folder)
    xml_files = list(xml_dir.glob("*.xml"))
    trackedTeam = str(args.team)
    sm = StatMonkey(xml_files, trackedTeam)

    stats1 = sm.getAllStats()
    sm.addAllStatsToPres(stats1)


if __name__ == "__main__":
    main()
