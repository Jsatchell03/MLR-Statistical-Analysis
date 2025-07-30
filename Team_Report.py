# Get all stats where a team is top 3 or bottom 3 then graph
# Grab league averages
# outlier stats array of objects with title: str, values: sorted arr(int), labels: arr(str)
# Function to uplad new file or just pull from data base. Include option to
# add XMLs
from pptx import Presentation
from pptx.util import Inches
import logging
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import argparse
import statistics
from pathlib import Path
import shutil
from collections import Counter


class TeamReport:
    ignore_list = [
        "Scrums Won",
        "Scrums Lost",
        "Total Restarts",
        "Restarts Retained",
        "Num Rucks Won",
        "Num Rucks Lost",
        "Mauls Won",
        "Mauls Lost",
    ]

    def __init__(self, team, file_path):
        self.prs = Presentation()
        self.outlier_stats = []
        self.stats_covered = []
        self.team = team
        self.excel_file = file_path
        self.all_sheets_data = pd.read_excel(
            self.excel_file,
            sheet_name=[
                "Teams Average",
                "Kicks",
                "Mauls",
                "Turnover Won",
                "Turnover Con",
                "Penalties",
                "Tackles",
                "Carries",
                "Ruck Entries",
                "Rucks",
                "Lineouts",
                "Scrums",
                "Restarts",
                "22m Entries",
                "Tries Overview",
                "Try Times",
            ],
            header=3,
        )

    def get_full_report(self, team):
        pass

    def get_outlier_stats(self):
        for sheet in self.all_sheets_data.keys():
            match sheet:
                case (
                    "Teams Average"
                    | "Lineouts"
                    | "Restarts"
                    | "22m Entries"
                    | "Tries Overview"
                ):
                    title_df = pd.read_excel(
                        self.excel_file,
                        sheet_name=sheet,
                        header=None,
                        skiprows=range(1, 15),
                        nrows=2,
                    )
                    sheet_title = title_df.iloc[0, 0]
                    sheet_title_opps = title_df.iloc[1, 0]
                    team_df = self.all_sheets_data[sheet].iloc[:11]
                    opp_df = self.all_sheets_data[sheet].iloc[15:27]
                    for col in team_df.columns[1:]:
                        if col in self.ignore_list or col in self.stats_covered:
                            continue
                        self.stats_covered.append(col)
                        sorted_team_df = team_df.sort_values(by=col, ascending=False)
                        sorted_opp_df = opp_df.sort_values(by=col, ascending=False)
                        sorted_teams = sorted_team_df.iloc[:, 0].tolist()
                        sorted_opps = sorted_opp_df.iloc[:, 0].tolist()
                        team_value = sorted_team_df[col].iloc[
                            sorted_teams.index(self.team)
                        ]
                        team_values = sorted_team_df[col].values
                        team_values_set = set(team_values)
                        opp_value = sorted_opp_df[col].iloc[
                            sorted_opps.index(self.team)
                        ]
                        opp_values = sorted_opp_df[col].values
                        opp_values_set = set(opp_values)
                        # If all values are the same
                        if len(team_values_set) <= 1 or len(opp_values_set) <= 1:
                            continue
                        # If there are 2 unique values and our team is not the outlier
                        elif (
                            len(team_values_set) <= 2
                            and team_value
                            == max(team_values, key=Counter(team_values).get)
                        ) or (
                            len(opp_values_set) <= 2
                            and opp_value
                            == max(opp_values, key=Counter(opp_values).get)
                        ):
                            continue
                        if (
                            self.team in sorted_teams[:3]
                            or self.team in sorted_teams[-3:]
                        ):
                            stat = {
                                "title": sheet_title + ": " + col,
                                "value": team_value,
                                "rank": sorted_teams.index(self.team) + 1,
                                "values": team_values,
                                "sorted_teams": sorted_teams,
                            }
                            self.outlier_stats.append(stat)

                        if (
                            self.team in sorted_opps[:3]
                            or self.team in sorted_opps[-3:]
                        ):
                            stat = {
                                "title": sheet_title_opps + " " + col,
                                "value": opp_value,
                                "rank": sorted_opps.index(self.team) + 1,
                                "values": opp_values,
                                "sorted_teams": sorted_opps,
                            }
                            self.outlier_stats.append(stat)
                case (
                    "Kicks"
                    | "Turnover Won"
                    | "Turnover Con"
                    | "Penalties"
                    | "Tackles"
                    | "Carries"
                    | "Ruck Entries"
                    | "Rucks"
                ):
                    title_df = pd.read_excel(
                        self.excel_file,
                        sheet_name=sheet,
                        header=None,
                        skiprows=15,
                        nrows=2,
                    )
                    sheet_title = title_df.iloc[0, 0]
                    df = self.all_sheets_data[sheet].iloc[15:27]
                    for col in df.columns[1:]:
                        if col in self.ignore_list or col in self.stats_covered:
                            continue
                        sorted_df = df.sort_values(by=col, ascending=False)
                        sorted_teams = sorted_df.iloc[:, 0].tolist()
                        values = sorted_df[col].values
                        values_set = set(values)
                        team_value = sorted_df[col].iloc[sorted_teams.index(self.team)]
                        # If all values are the same
                        if len(values_set) <= 1:
                            continue
                        # If there are 2 unique values and our team is not the outlier
                        elif len(values_set) <= 2 and team_value == max(
                            values, key=Counter(values).get
                        ):
                            continue

                        if (
                            self.team in sorted_teams[:3]
                            or self.team in sorted_teams[-3:]
                        ):
                            stat = {
                                "title": sheet_title + ": " + col,
                                "value": team_value,
                                "rank": sorted_teams.index(self.team) + 1,
                                "values": values,
                                "sorted_teams": sorted_teams,
                            }
                            self.outlier_stats.append(stat)

        return self.outlier_stats

    def draw_stats(self):
        if Path("graphs").is_dir():
            shutil.rmtree(Path("graphs"))
        Path("graphs").mkdir(parents=True, exist_ok=True)
        for stat in self.outlier_stats:
            path = "graphs/" + stat["title"]
            plt.figure(figsize=(11, 6))
            plt.title(stat["title"])
            colors = ["#1f77b4"] * len(stat["sorted_teams"])
            if self.team in stat["sorted_teams"]:
                team_index = stat["sorted_teams"].index(self.team)
                colors[team_index] = "#ff7f0e"
            bars = plt.bar(
                list(stat["sorted_teams"]), list(stat["values"]), color=colors
            )
            for bar in bars:
                height = bar.get_height()
                plt.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    (height / 2) if height != 0 else 1,
                    f"{round(height, 2)}",
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
            plt.xticks(rotation=45)
            plt.tight_layout()
            plt.subplots_adjust(bottom=0.25)
            plt.savefig(path)
            plt.close()

    def add_graphs_to_pres(self):
        graphs_dir = Path("graphs")
        for file_path in graphs_dir.iterdir():
            if file_path.is_file():
                self.add_stat_to_pres(str(file_path))

    def add_stat_to_pres(self, statImgPath):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        left = top = Inches(0)
        width = self.prs.slide_width
        height = self.prs.slide_height
        slide.shapes.add_picture("assets/bg.png", left, top, width, height)

        top = Inches(1.5)
        left = Inches(2.5)
        slide.shapes.add_picture(statImgPath, left, top)

        top = Inches(0)
        left = Inches(0)
        width = height = Inches(1.5)
        slide.shapes.add_picture(
            "assets/HoundsBadge_LightOnDarkBG.png", left, top, width, height
        )
        if self.team != "Chicago Hounds":
            top = Inches(0)
            left = Inches(2)
            width = height = Inches(1)
            slide.shapes.add_picture(
                f"assets/League Logos/{self.team.replace(' ', "_")}.png",
                left,
                top,
                width,
                height,
            )

        top = Inches(7.25)
        left = Inches(14.5)
        width = height = Inches(1.5)
        slide.shapes.add_picture(
            "assets/HoundsShield_LightOnDarkBG.png", left, top, width, height
        )

        self.prs.save(f"{self.team}_Full_Report.pptx")

    def check_outlier(self, stat, team_name):
        pass

    def build_comparison(self):
        pass


def main():
    parser = argparse.ArgumentParser(
        description="Generate a presentation with key stats for an MLR team."
    )
    parser.add_argument(
        "excel_file",
        help="File name for the Team Season Report provided by Oval Insights",
    )
    parser.add_argument(
        "team",
        help="Team name spelt and capitalize the exact way it is referenced in Oval Insights file",
    )
    args = parser.parse_args()
    tr = TeamReport(args.team, args.excel_file)
    # print(tr.get_outlier_stats())
    tr.get_outlier_stats()
    tr.draw_stats()
    tr.add_graphs_to_pres()


if __name__ == "__main__":
    main()
